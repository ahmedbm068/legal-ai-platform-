"""
Generate a professional PFE presentation PowerPoint for Legal AI Platform.
Run: .\.venv\Scripts\python.exe scripts/generate_presentation.py
Output: docs/Legal_AI_Platform_Soutenance.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ─────────────────────────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x0D, 0x1B, 0x2A)   # slide backgrounds
ACCENT_GOLD = RGBColor(0xC9, 0xA0, 0x3C)   # headings / highlights
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)
MID_GREY    = RGBColor(0x8A, 0x95, 0xA5)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
CARD_BG     = RGBColor(0x16, 0x2A, 0x44)
GREEN_OK    = RGBColor(0x2E, 0xCC, 0x71)
RED_RISK    = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # totally blank

ARCH_IMG  = os.path.join(os.path.dirname(__file__), "..", "docs", "mermaid", "images", "architecture_flow.png")
ORCH_IMG  = os.path.join(os.path.dirname(__file__), "..", "docs", "mermaid", "images", "ai_orchestration_flow.png")
ER_IMG    = os.path.join(os.path.dirname(__file__), "..", "docs", "mermaid", "images", "data_model_er.png")


# ── helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def fill_slide(slide, color=DARK_NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def txb(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return box

def rect(slide, l, t, w, h, fill_color=CARD_BG, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def gold_bar(slide, t=1.05, h=0.04):
    rect(slide, 0, t, 13.33, h, ACCENT_GOLD)

def section_title(slide, title, subtitle=""):
    fill_slide(slide)
    rect(slide, 0, 3.0, 13.33, 1.6, CARD_BG)
    gold_bar(slide, 2.95, 0.06)
    gold_bar(slide, 4.56, 0.06)
    txb(slide, title,    1.0, 3.05, 11.0, 1.0,  size=44, bold=True,  align=PP_ALIGN.CENTER)
    if subtitle:
        txb(slide, subtitle, 1.0, 4.1,  11.0, 0.5, size=20, color=MID_GREY, align=PP_ALIGN.CENTER)

def slide_header(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.0, CARD_BG)
    gold_bar(slide, 1.0, 0.05)
    txb(slide, title,    0.3, 0.08, 10.0, 0.55, size=26, bold=True)
    if subtitle:
        txb(slide, subtitle, 0.3, 0.62, 10.0, 0.35, size=14, color=ACCENT_GOLD)
    # slide number placeholder text bottom right
    return

def bullet_list(slide, items, l, t, w, h, size=16, color=WHITE, spacing=0.38):
    for i, item in enumerate(items):
        marker = "▸  " if not item.startswith("  ") else "      –  "
        txb(slide, marker + item.strip(), l, t + i * spacing, w, spacing + 0.05,
            size=size, color=color)

def img(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        # placeholder box
        r = rect(slide, l, t, w, h, CARD_BG, ACCENT_GOLD)
        txb(slide, "[image not found]", l+0.1, t+h/2-0.2, w-0.2, 0.4,
            size=12, color=MID_GREY, align=PP_ALIGN.CENTER)

def card(slide, title, body_lines, l, t, w, h,
         title_color=ACCENT_GOLD, body_size=14):
    rect(slide, l, t, w, h, CARD_BG, ACCENT_GOLD)
    txb(slide, title, l+0.12, t+0.1, w-0.24, 0.35,
        size=15, bold=True, color=title_color)
    y = t + 0.48
    for line in body_lines:
        txb(slide, line, l+0.15, y, w-0.3, 0.28, size=body_size, color=WHITE)
        y += 0.29

def connector_arrow(slide, x1, y1, x2, y2):
    """Simple horizontal or vertical connector drawn as a thin gold rect."""
    if abs(x2 - x1) >= abs(y2 - y1):
        rect(slide, min(x1,x2), (y1+y2)/2 - 0.02,
             abs(x2-x1), 0.04, ACCENT_GOLD)
    else:
        rect(slide, (x1+x2)/2 - 0.02, min(y1,y2),
             0.04, abs(y2-y1), ACCENT_GOLD)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
rect(sl, 0, 0, 13.33, 0.5, ACCENT_GOLD)
rect(sl, 0, 7.0, 13.33, 0.5, ACCENT_GOLD)
rect(sl, 0, 0.5, 13.33, 6.5, DARK_NAVY)
rect(sl, 0.4, 1.1, 12.53, 4.5, CARD_BG)

txb(sl, "Legal AI Platform",
    0.8, 1.3, 11.6, 1.2, size=52, bold=True, align=PP_ALIGN.CENTER, color=ACCENT_GOLD)
txb(sl, "AI-Powered Legal Workspace for Modern Law Firms",
    0.8, 2.55, 11.6, 0.65, size=22, align=PP_ALIGN.CENTER, color=WHITE)
txb(sl, "Case management  ·  Document Intelligence  ·  Multi-Agent Copilot  ·  Client Portal  ·  Voice Intake",
    0.8, 3.25, 11.6, 0.5, size=14, align=PP_ALIGN.CENTER, color=MID_GREY, italic=True)

gold_bar(sl, 3.85, 0.04)

txb(sl, "Presented by:   Ben Ahmed Mohamed",
    1.5, 4.2, 10.0, 0.4, size=16, align=PP_ALIGN.CENTER, color=WHITE)
txb(sl, "Company Supervisor :                University Supervisor :",
    1.5, 4.7, 10.0, 0.35, size=14, align=PP_ALIGN.CENTER, color=MID_GREY)
txb(sl, "Graduation Year : 2025-2026",
    1.5, 5.15, 10.0, 0.35, size=14, align=PP_ALIGN.CENTER, color=MID_GREY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — TABLE OF CONTENTS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Table of Contents")

items_l = [
    ("01", "Host Company Presentation"),
    ("02", "Introduction & Context"),
    ("03", "Study of Existing Solutions"),
    ("04", "Proposed Solution"),
    ("05", "Adopted Methodology"),
]
items_r = [
    ("06", "Project Analysis"),
    ("07", "Design & Technological Choices"),
    ("08", "Realization"),
    ("09", "Conclusion & Perspectives"),
    ("",   ""),
]

for i, (num, label) in enumerate(items_l):
    y = 1.4 + i * 0.95
    rect(sl, 0.5, y, 0.55, 0.65, ACCENT_GOLD)
    txb(sl, num, 0.5, y+0.1, 0.55, 0.45, size=18, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    txb(sl, label, 1.2, y+0.12, 5.0, 0.45, size=16, color=WHITE)

for i, (num, label) in enumerate(items_r):
    if not num:
        continue
    y = 1.4 + i * 0.95
    rect(sl, 7.0, y, 0.55, 0.65, ACCENT_GOLD)
    txb(sl, num, 7.0, y+0.1, 0.55, 0.45, size=18, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    txb(sl, label, 7.7, y+0.12, 5.0, 0.45, size=16, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — SECTION: HOST COMPANY
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_title(sl, "HOST COMPANY PRESENTATION", "01")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — HOST COMPANY DETAIL
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Host Company", "Where the project lives")

txb(sl, "⚖  Tunisian Law Firm — Modern Legal Practice",
    0.5, 1.25, 8.0, 0.55, size=20, bold=True, color=ACCENT_GOLD)
txb(sl, "The firm handles civil, commercial, and corporate legal matters across\nTunisia and cross-border jurisdictions including Germany.",
    0.5, 1.85, 7.5, 0.9, size=15, color=WHITE)

cards = [
    ("Case Management",    ["Multi-client case files", "Document tracking", "Deadline monitoring"]),
    ("Client Relations",   ["Intake workflows", "Status updates", "Portal access"]),
    ("Legal Research",     ["Code retrieval", "Jurisprudence lookup", "Drafting support"]),
    ("Compliance",         ["Multi-jurisdiction", "Audit trails", "Role isolation"]),
]
for i, (title, lines) in enumerate(cards):
    col = i % 2
    row = i // 2
    card(sl, title, lines, 0.45 + col * 6.25, 2.95 + row * 1.75, 5.9, 1.55)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — SECTION: INTRODUCTION
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_title(sl, "INTRODUCTION & CONTEXT", "02")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Problem Statement", "What lawyers face today")

# Before column
rect(sl, 0.4, 1.2, 5.7, 5.6, CARD_BG)
txb(sl, "BEFORE", 0.4, 1.2, 5.7, 0.5, size=16, bold=True, color=RED_RISK, align=PP_ALIGN.CENTER)
gold_bar(sl, 1.72, 0.03)
before = [
    "Documents scattered across tools",
    "Manual case summaries (hours)",
    "No jurisdiction-aware reasoning",
    "Generic AI with no case context",
    "No client self-service portal",
    "Voice intake requires manual entry",
    "No grounded citations or audit trail",
]
for i, line in enumerate(before):
    txb(sl, "✗  " + line, 0.6, 1.85 + i*0.6, 5.3, 0.5, size=14, color=WHITE)

# After column
rect(sl, 6.8, 1.2, 6.1, 5.6, CARD_BG)
txb(sl, "AFTER", 6.8, 1.2, 6.1, 0.5, size=16, bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
gold_bar(sl, 1.72, 0.03)
after = [
    "Unified case workspace with search",
    "AI copilot: instant grounded summaries",
    "Tunisia + Germany legal code retrieval",
    "Case-scoped multi-agent reasoning",
    "Separate secure client portal",
    "Voice → transcript → consultation",
    "Source citations + confidence scores",
]
for i, line in enumerate(after):
    txb(sl, "✓  " + line, 7.0, 1.85 + i*0.6, 5.7, 0.5, size=14, color=WHITE)

# Arrow between columns
rect(sl, 6.14, 3.75, 0.62, 0.12, ACCENT_GOLD)
txb(sl, "➜", 6.14, 3.6, 0.62, 0.5, size=24, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — SECTION: EXISTING SOLUTIONS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_title(sl, "STUDY OF EXISTING SOLUTIONS", "03")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — COMPETITIVE ANALYSIS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Competitive Analysis", "Where existing tools fall short")

cols = ["Solution", "Strengths ✓", "Limitations ✗"]
col_x = [0.3, 3.5, 8.5]
col_w = [3.0, 4.8, 4.6]

# header row
for ci, (c, x, w) in enumerate(zip(cols, col_x, col_w)):
    rect(sl, x, 1.2, w, 0.45, ACCENT_GOLD)
    txb(sl, c, x+0.08, 1.25, w-0.16, 0.38, size=14, bold=True, color=DARK_NAVY)

rows = [
    ("Harvey AI",       "LLM for lawyers · Citation focus · Trusted by large firms",
                        "US-centric · No Tunisia/MENA · Proprietary · No custom corpus"),
    ("Westlaw / LexisNexis",
                        "Huge corpus · Reliable precedents",
                        "Search-only · No generative AI copilot · Very expensive"),
    ("CaseText / Clio", "Practice management · Some AI features",
                        "Generic AI, not grounded · No multi-agent pipeline"),
    ("ChatGPT / Gemini","Fast generation · Good language",
                        "No case scoping · Hallucinations · No legal citations"),
    ("Our Platform",    "Case-grounded RAG · Multi-agent · Voice · Client portal · Tunisia legal codes",
                        "—  Purpose-built for this firm"),
]

for ri, (sol, strengths, limits) in enumerate(rows):
    y = 1.7 + ri * 1.0
    bg = CARD_BG if ri < 4 else RGBColor(0x0A, 0x3D, 0x2B)
    for ci, (x, w) in enumerate(zip(col_x, col_w)):
        rect(sl, x, y, w, 0.88, bg)
    txb(sl, sol,      col_x[0]+0.08, y+0.08, col_w[0]-0.16, 0.75, size=13, bold=(ri==4), color=ACCENT_GOLD if ri==4 else WHITE)
    txb(sl, strengths, col_x[1]+0.08, y+0.08, col_w[1]-0.16, 0.75, size=12, color=GREEN_OK if ri==4 else WHITE)
    txb(sl, limits,   col_x[2]+0.08, y+0.08, col_w[2]-0.16, 0.75, size=12, color=MID_GREY if ri==4 else RGBColor(0xFF,0x99,0x80))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — SECTION: PROPOSED SOLUTION
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_title(sl, "PROPOSED SOLUTION", "04")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — SOLUTION OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Proposed Solution", "A complete legal intelligence platform")

pillars = [
    ("⚖", "Case\nWorkspace",   "Cases · Clients · Documents\nDeadlines · RBAC · Multi-tenant"),
    ("🤖", "AI Copilot",       "Multi-agent orchestration\nIntent routing · Grounded answers"),
    ("📄", "Document\nIntel",  "OCR · NER · PII redaction\nChunking · FAISS vector index"),
    ("🔍", "Legal Search",     "Tunisia & Germany legal codes\nCitation verification · Fallback taxonomy"),
    ("🎤", "Voice Intake",     "Transcription pipeline\nAudio → Consultation → Case"),
    ("🌐", "Client\nPortal",   "Public intake form\nCase status tracking · Secure login"),
]

for i, (icon, title, desc) in enumerate(pillars):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.28
    y = 1.4 + row * 2.65
    rect(sl, x, y, 3.95, 2.4, CARD_BG, ACCENT_GOLD)
    txb(sl, icon,  x+0.15, y+0.12, 0.6, 0.55, size=26)
    txb(sl, title, x+0.8,  y+0.12, 2.9, 0.6, size=15, bold=True, color=ACCENT_GOLD)
    txb(sl, desc,  x+0.15, y+0.82, 3.6, 1.4, size=12, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — SECTION: METHODOLOGY
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_title(sl, "ADOPTED METHODOLOGY", "05")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — METHODOLOGY
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Methodology", "Scrum + iterative AI development")

# Scrum cycle visual
phases = ["Sprint\nPlanning", "Daily\nStandup", "Sprint\nExecution", "Sprint\nReview", "Retrospective"]
cx = [1.2, 3.5, 6.6, 9.7, 11.9]
for i, (phase, x) in enumerate(zip(phases, cx)):
    rect(sl, x, 1.35, 1.7, 1.3, CARD_BG, ACCENT_GOLD)
    txb(sl, phase, x+0.08, 1.5, 1.55, 1.0, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(phases)-1:
        txb(sl, "→", cx[i]+1.72, 1.8, 0.5, 0.5, size=20, bold=True, color=ACCENT_GOLD)

# Sprint table
sprint_data = [
    ("Sprint 1–2",  "Feb 2026",  "Backend core: auth, cases, clients, documents, RBAC, multi-tenant"),
    ("Sprint 3–4",  "Mar 2026",  "Document AI pipeline: OCR, NER, PII, chunking, FAISS indexing, RAG service"),
    ("Sprint 5–6",  "Mar 2026",  "Copilot orchestration, intent routing, specialized agents, voice intake"),
    ("Sprint 7–8",  "Apr 2026",  "Internal React workspace: case view, document panel, copilot chat, calendar"),
    ("Sprint 9–10", "Apr 2026",  "Client portal: intake form, status tracking, portal auth, MinIO uploads"),
    ("Sprint 11–12","May 2026",  "Trust hardening: grounding, citation verification, AI Insight panel, eval suite"),
]

rect(sl, 0.3, 2.85, 12.73, 0.4, ACCENT_GOLD)
txb(sl, "Sprint", 0.35, 2.88, 1.4, 0.35, size=13, bold=True, color=DARK_NAVY)
txb(sl, "Period", 1.85, 2.88, 1.3, 0.35, size=13, bold=True, color=DARK_NAVY)
txb(sl, "Deliverables", 3.25, 2.88, 9.6, 0.35, size=13, bold=True, color=DARK_NAVY)

for i, (sp, period, deliv) in enumerate(sprint_data):
    y = 3.3 + i * 0.63
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x11, 0x22, 0x38)
    rect(sl, 0.3, y, 12.73, 0.58, bg)
    txb(sl, sp,     0.38, y+0.06, 1.35, 0.45, size=12, bold=True, color=ACCENT_GOLD)
    txb(sl, period, 1.85, y+0.06, 1.25, 0.45, size=12, color=MID_GREY)
    txb(sl, deliv,  3.25, y+0.06, 9.6,  0.45, size=12, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — SECTION: PROJECT ANALYSIS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_title(sl, "PROJECT ANALYSIS", "06")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — PROJECT ACTORS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Project Actors", "Who uses the platform and how")

actors = [
    ("🧑‍⚖️", "Lawyer",      ["Primary user", "Manages cases & clients", "Uses copilot & legal search", "Reviews AI-grounded answers", "Signs off on documents"]),
    ("🛡️", "Admin",        ["Manages user accounts", "Tenant configuration", "Monitors system usage", "Access control & RBAC"]),
    ("🤝", "Legal\nAssistant", ["Uploads documents", "Tracks deadlines", "Runs document summaries", "Prepares drafts"]),
    ("👤", "Client",       ["Submits intake form", "Tracks case status", "Uploads supporting docs", "Receives status updates"]),
]

for i, (icon, name, roles) in enumerate(actors):
    x = 0.3 + i * 3.22
    rect(sl, x, 1.3, 3.0, 5.5, CARD_BG, ACCENT_GOLD)
    txb(sl, icon, x+1.1, 1.4, 0.9, 0.7, size=30, align=PP_ALIGN.CENTER)
    txb(sl, name, x+0.1, 2.18, 2.8, 0.55, size=16, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    gold_bar(sl, 2.78, 0.03)
    for j, role in enumerate(roles):
        txb(sl, "▸ " + role, x+0.18, 2.9 + j*0.62, 2.65, 0.5, size=12, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — FUNCTIONAL REQUIREMENTS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Functional Requirements")

groups = [
    ("Case & Client Management",
     ["Create/update/close legal cases", "Manage clients and contacts", "Track deadlines and obligations", "Multi-tenant data isolation"]),
    ("AI Copilot & Agents",
     ["Intent-routed copilot pipeline", "Ask, summarize, draft, analyze intents", "Specialized agents per legal task", "Grounded answers with source citations"]),
    ("Document Intelligence",
     ["Upload PDF, DOCX, scanned images", "OCR, NER, PII redaction pipeline", "Chunking, embedding, FAISS indexing", "Hybrid retrieval with reranking"]),
    ("Legal Search",
     ["Tunisia + Germany legal codes corpus", "Article-level citation retrieval", "Claim verification agent", "Fallback with explicit reason taxonomy"]),
    ("Voice & Intake",
     ["Voice recording upload / live capture", "Transcription to structured text", "Auto-extract consultation fields", "Consultation linked to case"]),
    ("Client Portal",
     ["Public intake form submission", "Case status self-service tracking", "Secure portal login (OTP)", "Document upload by client"]),
]

for i, (title, items) in enumerate(groups):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.55
    y = 1.25 + row * 2.05
    rect(sl, x, y, 6.25, 1.9, CARD_BG, ACCENT_GOLD)
    txb(sl, title, x+0.12, y+0.1, 5.9, 0.38, size=14, bold=True, color=ACCENT_GOLD)
    for j, item in enumerate(items):
        txb(sl, "▸ " + item, x+0.18, y+0.55 + j*0.32, 5.85, 0.3, size=11.5, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — NON-FUNCTIONAL REQUIREMENTS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Non-Functional Requirements")

nfr = [
    ("🔒 Security",      ["JWT + refresh token auth", "BCrypt password hashing", "RBAC with role isolation", "Multi-tenant data boundaries", "CORS + rate limiting"]),
    ("⚡ Performance",   ["Background workers for heavy AI tasks", "Redis caching for frequent lookups", "FAISS in-memory vector search", "Paginated API endpoints", "Lazy-loaded frontend bundles"]),
    ("🛡️ Reliability",  ["LLM provider fallback chain", "Structured fallback reason taxonomy", "Graceful degradation on 429 errors", "Smoke tests + E2E regression suite", "Eval gates before release"]),
    ("🏗️ Maintainability", ["Modular service architecture", "Typed API contracts (Pydantic)", "Clear agent-per-intent separation", "Centralized prompt contracts", "Documented sprint history"]),
]

for i, (title, items) in enumerate(nfr):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.55
    y = 1.25 + row * 2.8
    rect(sl, x, y, 6.25, 2.6, CARD_BG, ACCENT_GOLD)
    txb(sl, title, x+0.12, y+0.1, 5.9, 0.42, size=15, bold=True, color=ACCENT_GOLD)
    for j, item in enumerate(items):
        txb(sl, "▸ " + item, x+0.18, y+0.62 + j*0.38, 5.85, 0.35, size=12.5, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — USE CASE DIAGRAM (drawn with shapes)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Use Case Diagram", "Actor–feature mapping")

# System boundary box
rect(sl, 2.2, 1.15, 9.0, 5.8, RGBColor(0x10, 0x20, 0x35), ACCENT_GOLD)
txb(sl, "Legal AI Platform System", 2.2, 1.15, 9.0, 0.4, size=12, italic=True,
    color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# Actors on left
actors_uc = [("Lawyer", 1.05, 2.2), ("Admin", 1.05, 4.0), ("Client", 1.05, 5.8)]
for name, ax, ay in actors_uc:
    # stick figure approximation: circle + box
    shape = sl.shapes.add_shape(9, Inches(ax+0.25), Inches(ay-0.28), Inches(0.45), Inches(0.45))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT_GOLD
    shape.line.fill.background()
    txb(sl, name, ax, ay+0.22, 0.95, 0.3, size=11, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# Use cases per actor
uc_rows = [
    # (label, x, y, actors connected [left_x, left_y list])
    ("Manage Cases & Clients",         3.1, 1.8),
    ("Ask AI Copilot (case-scoped)",   3.1, 2.6),
    ("Upload & Search Documents",      3.1, 3.4),
    ("Legal Code Search",              3.1, 4.2),
    ("Voice Intake → Consultation",    3.1, 5.0),
    ("Manage Users & Tenants",         7.2, 2.4),
    ("View System Audit Logs",         7.2, 3.4),
    ("Submit Intake Form",             7.2, 4.4),
    ("Track Case Status",              7.2, 5.4),
]
for label, ux, uy in uc_rows:
    # Ellipse
    shape = sl.shapes.add_shape(9, Inches(ux), Inches(uy-0.22),
                                Inches(3.6), Inches(0.48))
    shape.fill.solid(); shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = ACCENT_GOLD
    txb(sl, label, ux+0.08, uy-0.18, 3.44, 0.38, size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — SEQUENCE DIAGRAM (copilot pipeline)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Sequence Diagram — Copilot Pipeline", "From user prompt to grounded answer")

# Lifeline headers
lifelines = ["Lawyer UI", "API Gateway", "Orchestrator", "RAG Service", "Verifier Agent", "LLM Gateway", "Assembly"]
lx = [0.5, 2.1, 3.75, 5.4, 7.05, 8.7, 10.5]
lw = 1.3

for i, (label, x) in enumerate(zip(lifelines, lx)):
    rect(sl, x, 1.15, lw, 0.5, ACCENT_GOLD if i in (0, 6) else CARD_BG, WHITE)
    txb(sl, label, x+0.05, 1.18, lw-0.1, 0.42, size=10, bold=True,
        color=DARK_NAVY if i in (0,6) else WHITE, align=PP_ALIGN.CENTER)
    # Dashed lifeline
    for seg in range(12):
        rect(sl, x + lw/2 - 0.02, 1.7 + seg*0.42, 0.04, 0.3, ACCENT_GOLD)

# Messages
messages = [
    (0, 1, 1.9,  "POST /copilot  {prompt, case_id, scope}"),
    (1, 2, 2.35, "dispatch(intent, context)"),
    (2, 3, 2.8,  "hybrid_retrieve(query, case_id)"),
    (3, 4, 3.25, "verify_sources(chunks)"),
    (4, 3, 3.7,  "verified_chunks + confidence"),
    (3, 2, 4.15, "grounded_context"),
    (2, 5, 4.6,  "generate(prompt + context)"),
    (5, 2, 5.05, "raw_answer"),
    (2, 6, 5.5,  "assemble(answer, grounding, ai_insight)"),
    (6, 1, 5.95, "structured_response"),
    (1, 0, 6.4,  "200 OK  {answer, confidence, sources, ai_insight}"),
]

for src, dst, my, label in messages:
    sx = lx[src] + lw/2
    dx = lx[dst] + lw/2
    # line
    rect(sl, min(sx,dx), my, abs(dx-sx), 0.035, ACCENT_GOLD)
    # arrowhead stub
    arrow_x = dx - 0.12 if dx > sx else dx
    rect(sl, arrow_x, my-0.06, 0.14, 0.1, ACCENT_GOLD)
    # label
    mid_x = (sx + dx) / 2 - 0.8
    txb(sl, label, mid_x, my-0.26, 1.8, 0.22, size=8.5, color=MID_GREY)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — CLASS / ENTITY DIAGRAM (ER diagram image)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Entity-Relationship Diagram", "Core data model")

img(sl, ER_IMG, 0.3, 1.15, 12.7, 6.1)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — SECTION: DESIGN & TECH CHOICES
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_title(sl, "DESIGN & TECHNOLOGICAL CHOICES", "07")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 21 — LOGICAL ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Logical Architecture", "Layers and responsibilities")

img(sl, ARCH_IMG, 0.3, 1.15, 12.7, 6.1)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 22 — AI ORCHESTRATION FLOW
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "AI Orchestration Flow", "Multi-agent pipeline internals")

img(sl, ORCH_IMG, 0.3, 1.15, 12.7, 6.1)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 23 — PHYSICAL ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Physical Architecture", "Docker deployment topology")

# Draw the Docker Compose stack visually

# Host machine box
rect(sl, 0.25, 1.1, 12.8, 6.15, RGBColor(0x08, 0x18, 0x28), ACCENT_GOLD)
txb(sl, "Docker Host", 0.3, 1.1, 3.0, 0.4, size=12, italic=True, color=ACCENT_GOLD)

# Frontend containers
rect(sl, 0.5, 1.6, 2.8, 1.2, CARD_BG, WHITE)
txb(sl, "React Frontend\n(Internal Workspace)", 0.6, 1.65, 2.6, 1.05, size=12, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "Port 5173", 0.6, 2.55, 2.6, 0.25, size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

rect(sl, 3.55, 1.6, 2.8, 1.2, CARD_BG, WHITE)
txb(sl, "React Client Portal\n(Public Intake)", 3.65, 1.65, 2.6, 1.05, size=12, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "Port 5174", 3.65, 2.55, 2.6, 0.25, size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

# Backend
rect(sl, 6.6, 1.6, 3.0, 1.2, RGBColor(0x0A, 0x35, 0x5A), ACCENT_GOLD)
txb(sl, "FastAPI Backend\n(Uvicorn --reload)", 6.7, 1.65, 2.8, 1.05, size=12, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
txb(sl, "Port 8000", 6.7, 2.55, 2.8, 0.25, size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

# External AI
rect(sl, 9.85, 1.6, 2.9, 1.2, RGBColor(0x2C, 0x1A, 0x52), WHITE)
txb(sl, "LLM Providers\n(Groq / OpenAI / OpenRouter)", 9.95, 1.65, 2.7, 1.05, size=11, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "External API", 9.95, 2.55, 2.7, 0.25, size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

# Infrastructure row
infra = [
    ("PostgreSQL 15\nPort 5433", 0.5),
    ("Redis 7\nPort 6379", 3.25),
    ("MinIO\nObject Storage\nPort 9000", 6.0),
    ("FAISS\nVector Index\n(In-memory)", 9.3),
]
for label, x in infra:
    rect(sl, x, 4.1, 2.6, 1.5, RGBColor(0x18, 0x30, 0x48), ACCENT_GOLD)
    txb(sl, label, x+0.1, 4.2, 2.4, 1.25, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# AI Services row
rect(sl, 0.5, 5.85, 12.2, 0.9, RGBColor(0x0C, 0x28, 0x40), ACCENT_GOLD)
ai_services = ["Document AI Pipeline", "RAG Service", "Legal Search", "Copilot Orchestrator",
               "Voice / Whisper", "Specialized Agents x10"]
for i, svc in enumerate(ai_services):
    txb(sl, "▸ " + svc, 0.65 + i * 2.0, 5.95, 1.92, 0.35, size=10, color=ACCENT_GOLD)

txb(sl, "AI Services Layer", 0.6, 5.88, 2.5, 0.3, size=10, bold=True, italic=True, color=MID_GREY)

# Arrows: frontends → backend
for fx in [1.9, 4.95]:
    rect(sl, fx, 2.82, 0.04, 1.27, ACCENT_GOLD)
    txb(sl, "↓", fx-0.1, 3.9, 0.3, 0.3, size=14, color=ACCENT_GOLD)

# Backend → infra
for ix in [1.8, 4.55, 7.3, 10.6]:
    rect(sl, ix, 2.82, 0.04, 1.27, ACCENT_GOLD)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 24 — SOFTWARE STACK
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Software Stack", "Technologies chosen and why")

layers = [
    ("Frontend",    ACCENT_GOLD,   "React 18 + TypeScript + Vite + TailwindCSS",         "Two separate SPAs: internal workspace and public portal"),
    ("Backend",     RGBColor(0x21,0x8A,0xFF), "FastAPI + Python 3 + Pydantic + Uvicorn", "Async-first, auto-documented REST API with typed contracts"),
    ("Database",    RGBColor(0x00,0x8B,0x8B), "PostgreSQL 15 + SQLAlchemy ORM",          "Relational model with multi-tenant isolation per schema"),
    ("AI / ML",     RGBColor(0x9B,0x59,0xB6), "sentence-transformers + FAISS + Groq / OpenAI", "Hybrid retrieval: dense embeddings + BM25 rerank + LLM generation"),
    ("Storage",     RGBColor(0xE6,0x7E,0x22), "MinIO object storage",                    "Document binary storage (PDFs, audio) with presigned URLs"),
    ("Infra",       RGBColor(0x27,0xAE,0x60), "Docker Compose + Redis 7 + GitHub Actions", "Containerized local dev; CI trust-release-gate workflow"),
]

for i, (layer, color, tech, reason) in enumerate(layers):
    y = 1.25 + i * 0.99
    rect(sl, 0.3, y, 1.5, 0.82, color)
    txb(sl, layer, 0.3, y+0.18, 1.5, 0.45, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, 1.85, y, 5.0, 0.82, CARD_BG)
    txb(sl, tech, 2.0, y+0.12, 4.7, 0.6, size=12, color=WHITE)
    rect(sl, 6.9, y, 6.1, 0.82, RGBColor(0x10, 0x20, 0x33))
    txb(sl, reason, 7.05, y+0.12, 5.85, 0.6, size=12, color=MID_GREY)

txb(sl, "Layer", 0.35, 1.08, 1.4, 0.3, size=11, bold=True, color=ACCENT_GOLD)
txb(sl, "Technology", 2.0, 1.08, 4.7, 0.3, size=11, bold=True, color=ACCENT_GOLD)
txb(sl, "Design Rationale", 7.05, 1.08, 5.5, 0.3, size=11, bold=True, color=ACCENT_GOLD)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 25 — SECURITY DESIGN
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Security Design", "Defence-in-depth across all layers")

sec_items = [
    ("Authentication",
     ["JWT bearer tokens (stateless)", "BCrypt password hashing (cost factor 12)", "Refresh token rotation", "OTP login for client portal"]),
    ("Authorization",
     ["RBAC: Admin / Lawyer / Assistant / Client", "Multi-tenant: requests filtered by tenant_id", "Entity-level ownership checks on all routes", "Client portal isolated from internal API"]),
    ("Transport & Network",
     ["HTTPS enforced in production", "CORS allowlist per origin", "CORS preflight cache headers", "Rate limiting middleware"]),
    ("Data Protection",
     ["PII redaction before embedding", "MinIO presigned URLs (time-limited)", "No secrets in source code (.env pattern)", "Docker secrets for credentials"]),
]

for i, (title, points) in enumerate(sec_items):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.55
    y = 1.25 + row * 2.85
    rect(sl, x, y, 6.25, 2.65, CARD_BG, ACCENT_GOLD)
    # gold title bar
    rect(sl, x, y, 6.25, 0.45, ACCENT_GOLD)
    txb(sl, title, x+0.12, y+0.06, 5.9, 0.35, size=14, bold=True, color=DARK_NAVY)
    for j, pt in enumerate(points):
        txb(sl, "▸ " + pt, x+0.18, y+0.58 + j*0.5, 5.85, 0.44, size=12.5, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 26 — SECTION: REALIZATION
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_title(sl, "REALIZATION", "08")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 27 — REALIZATION: AI PIPELINE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Realization — AI Copilot Pipeline", "Sprint 5–6: Orchestration & agents")

# Pipeline flow boxes
steps = [
    ("1\nUser\nPrompt",         "Lawyer types\nnatural language"),
    ("2\nIntent\nDetection",    "15+ intents:\nask, draft, summarize\ntimeline, risks..."),
    ("3\nContext\nEnrichment",  "Case history,\nactive docs,\nprior answers"),
    ("4\nRAG\nRetrieval",       "FAISS dense +\nBM25 rerank\nhybrid search"),
    ("5\nLLM\nGeneration",      "Groq / OpenAI\nwith grounded\nprompt contracts"),
    ("6\nVerifier\nAgent",      "Citation check\nFallback reason\nConfidence score"),
    ("7\nAssembly\n& Insight",  "Structured answer\nAI Insight panel\nLawyer note"),
]

for i, (title, desc) in enumerate(steps):
    x = 0.35 + i * 1.84
    rect(sl, x, 1.3, 1.62, 1.85, CARD_BG, ACCENT_GOLD)
    txb(sl, title, x+0.08, 1.35, 1.46, 0.9, size=11, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    txb(sl, desc,  x+0.08, 2.28, 1.46, 0.82, size=9.5, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        txb(sl, "→", x+1.65, 1.95, 0.22, 0.4, size=16, bold=True, color=ACCENT_GOLD)

# Agents box
rect(sl, 0.35, 3.4, 12.63, 3.9, RGBColor(0x0C, 0x22, 0x38), ACCENT_GOLD)
txb(sl, "Specialized Agent Layer", 0.5, 3.42, 5.0, 0.38, size=13, bold=True, color=ACCENT_GOLD)

agents = [
    "PromptOptimizerAgent", "RetrievalAgent", "CaseReasoningAgent",
    "TimelineAgent", "DraftingAgent", "DocumentComparisonAgent",
    "StrictVerifierAgent", "ClaimValidationAgent", "ContradictionDetectionAgent",
    "MatterClassificationAgent", "ArticleApplicabilityAgent", "WorkflowOrchestrator",
]
for i, agent in enumerate(agents):
    col = i % 4
    row = i // 4
    x = 0.55 + col * 3.15
    y = 3.9 + row * 1.05
    rect(sl, x, y, 2.88, 0.88, CARD_BG, RGBColor(0x44, 0x55, 0x77))
    txb(sl, agent, x+0.1, y+0.18, 2.7, 0.5, size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 28 — REALIZATION: DOCUMENT PIPELINE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Realization — Document Intelligence Pipeline", "Sprint 3–4: From raw file to searchable knowledge")

pipe_steps = [
    ("📤 Upload",      "PDF / DOCX / Image\nMinIO storage\nMetadata in DB"),
    ("🔍 OCR",         "Tesseract / PyPDF\nScanned image support\nText cleaning"),
    ("🏷️ NER",         "Named entity\nrecognition\nPerson, Org, Date"),
    ("🔒 PII\nRedact",  "Remove sensitive\npersonal data\nbefore embedding"),
    ("✂️ Chunk",        "Sliding window\nchunking\nOverlap strategy"),
    ("🧮 Embed",        "sentence-transformers\nall-MiniLM-L6-v2\n384-dim vectors"),
    ("📦 Index",        "FAISS CPU index\nPersisted to disk\nMetadata JSON"),
]

for i, (title, desc) in enumerate(pipe_steps):
    x = 0.35 + i * 1.84
    rect(sl, x, 1.3, 1.62, 2.0, CARD_BG, ACCENT_GOLD)
    txb(sl, title, x+0.08, 1.36, 1.46, 0.68, size=12, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    txb(sl, desc,  x+0.08, 2.1,  1.46, 1.1,  size=9.5, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(pipe_steps)-1:
        txb(sl, "→", x+1.65, 2.1, 0.22, 0.4, size=16, bold=True, color=ACCENT_GOLD)

# Retrieval section
rect(sl, 0.35, 3.55, 12.63, 3.7, RGBColor(0x0C, 0x22, 0x38), ACCENT_GOLD)
txb(sl, "Hybrid Retrieval at Query Time", 0.5, 3.58, 6.0, 0.38, size=13, bold=True, color=ACCENT_GOLD)

retrieval_steps = [
    ("Query\nEmbedding",    "Encode user query\nwith same model"),
    ("FAISS\nDense Search", "Top-K by\ncosine similarity"),
    ("BM25\nKeyword",       "Exact term match\nover corpus"),
    ("Cross-Encoder\nRerank","Score all candidates\ntogether"),
    ("Top-N\nChunks",       "Best N chunks\npass to LLM"),
    ("Source\nCitations",   "Article + page\nmetadata attached"),
]

for i, (title, desc) in enumerate(retrieval_steps):
    x = 0.55 + i * 2.1
    rect(sl, x, 4.1, 1.95, 2.9, CARD_BG, RGBColor(0x33, 0x66, 0x99))
    txb(sl, title, x+0.08, 4.17, 1.82, 0.6, size=11, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    txb(sl, desc,  x+0.08, 4.85, 1.82, 1.9, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(retrieval_steps)-1:
        txb(sl, "→", x+1.98, 5.4, 0.18, 0.35, size=14, bold=True, color=ACCENT_GOLD)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 29 — REALIZATION: AI INSIGHT & GROUNDING
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Realization — AI Insight & Trust Layer", "Grounding classification + confidence scoring")

txb(sl, "Every AI answer is accompanied by a structured AI Insight block:",
    0.4, 1.15, 12.0, 0.4, size=14, color=WHITE)

# AI Insight card mockup
rect(sl, 0.4, 1.65, 5.8, 5.5, RGBColor(0x10, 0x28, 0x44), ACCENT_GOLD)
txb(sl, "AI Insight  —  Example Response", 0.52, 1.72, 5.55, 0.38, size=13, bold=True, color=ACCENT_GOLD)

fields = [
    ("grounding_type",      "Case-grounded (document-based)"),
    ("confidence_level",    "high"),
    ("sources_count",       "7 document chunks matched"),
    ("legal_grounding",     "None — no verified legal provisions"),
    ("grounding_desc",      "Strongly grounded in case documents,\nnot supported by legal corpus"),
    ("lawyer_note",         "Case record supports this response.\nCounsel review advised before external use."),
]
for i, (key, val) in enumerate(fields):
    y = 2.2 + i * 0.82
    txb(sl, key + ":", 0.58, y, 2.2, 0.35, size=11, color=MID_GREY, italic=True)
    txb(sl, val,       2.85, y, 3.1, 0.65, size=11, color=GREEN_OK if "high" in val else WHITE)

# Grounding matrix
rect(sl, 6.5, 1.65, 6.55, 5.5, CARD_BG, ACCENT_GOLD)
txb(sl, "Grounding Classification Matrix", 6.62, 1.72, 6.3, 0.38, size=13, bold=True, color=ACCENT_GOLD)

matrix = [
    ("Scenario",                     "Grounding Type",                "Confidence"),
    ("3+ doc sources, no fallback",  "Case-grounded (document-based)","high"),
    ("1–2 doc sources",              "Case-grounded",                 "medium"),
    ("Case context, no legal codes", "Case-context reasoning",        "low"),
    ("External search fallback",     "Fallback (context-based)",      "medium"),
    ("No sources found",             "Not grounded",                  "low"),
    ("General chat (no case)",       "General",                       "varies"),
]
for i, (scen, gtype, conf) in enumerate(matrix):
    y = 2.22 + i * 0.72
    bg = ACCENT_GOLD if i == 0 else (RGBColor(0x0A,0x3D,0x2B) if conf == "high" else CARD_BG)
    rect(sl, 6.55, y, 6.4, 0.65, bg)
    c = DARK_NAVY if i == 0 else WHITE
    txb(sl, scen,  6.62, y+0.08, 2.5, 0.5, size=9.5, color=c, bold=(i==0))
    txb(sl, gtype, 9.2,  y+0.08, 2.3, 0.5, size=9.5, color=c, bold=(i==0))
    conf_color = GREEN_OK if conf == "high" else (RED_RISK if conf == "low" else MID_GREY) if i > 0 else DARK_NAVY
    txb(sl, conf, 11.55, y+0.08, 1.3, 0.5, size=9.5, color=conf_color, bold=(i==0))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 30 — REALIZATION: FRONTEND
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Realization — Frontend Workspaces", "Internal workspace + public client portal")

# Internal workspace
rect(sl, 0.3, 1.2, 6.1, 5.95, CARD_BG, ACCENT_GOLD)
txb(sl, "Internal Legal Workspace", 0.42, 1.25, 5.8, 0.45, size=14, bold=True, color=ACCENT_GOLD)
internal_pages = [
    ("Cases Page",       "Case list, filters, status badges, quick open"),
    ("Documents Page",   "Upload, preview, processing status, search"),
    ("Copilot Chat",     "Intent-aware chat, AI Insight panel, sources"),
    ("Calendar Page",    "Deadline view, sync, event management"),
    ("Evidence Review",  "Side-by-side doc comparison, annotations"),
    ("Admin Panel",      "User management, tenant config, audit logs"),
    ("Voice Intake",     "Record/upload audio, transcript view, link to case"),
]
for i, (page, desc) in enumerate(internal_pages):
    y = 1.82 + i * 0.72
    rect(sl, 0.42, y, 5.85, 0.62, RGBColor(0x10, 0x22, 0x38))
    txb(sl, page, 0.55, y+0.08, 1.9, 0.45, size=11, bold=True, color=ACCENT_GOLD)
    txb(sl, desc, 2.5,  y+0.08, 3.6, 0.45, size=10.5, color=WHITE)

# Client portal
rect(sl, 6.65, 1.2, 6.4, 5.95, CARD_BG, ACCENT_GOLD)
txb(sl, "Client Portal (Public)", 6.77, 1.25, 6.1, 0.45, size=14, bold=True, color=ACCENT_GOLD)
portal_pages = [
    ("Intake Form",      "Name, contact, matter type, free-text description"),
    ("Case Status",      "Track status: pending → in-progress → closed"),
    ("Document Upload",  "Attach supporting documents to intake request"),
    ("Portal Login",     "OTP-based login, isolated from staff auth"),
    ("Status History",   "Timeline of case updates and lawyer actions"),
]
for i, (page, desc) in enumerate(portal_pages):
    y = 1.82 + i * 0.86
    rect(sl, 6.77, y, 6.15, 0.76, RGBColor(0x10, 0x22, 0x38))
    txb(sl, page, 6.9, y+0.1, 2.0, 0.55, size=12, bold=True, color=ACCENT_GOLD)
    txb(sl, desc, 8.95, y+0.1, 3.75, 0.55, size=11, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 31 — REALIZATION: TESTS & QUALITY
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Realization — Tests & Quality Gates", "Automated validation across all trust-critical paths")

# Eval results
rect(sl, 0.3, 1.15, 12.73, 0.5, ACCENT_GOLD)
txb(sl, "Automated Eval Results", 0.42, 1.2, 6.0, 0.38, size=13, bold=True, color=DARK_NAVY)

eval_rows = [
    ("Full agent eval suite",          "43 / 43",  "100.0%", "advancement/evals/agent_eval_report_20260330.json"),
    ("Summary eval subset",            " 6 /  6",  "100.0%", "advancement/evals/agent_eval_report_summary.json"),
    ("E2E copilot smoke test",         " 4 /  4",  "100.0%", "scripts/test_copilot_e2e.py"),
    ("Backend unit tests (pytest)",    "All pass", "—",      "tests/"),
    ("Frontend build",                 "0 errors", "—",      "npm run build --prefix frontend"),
    ("Client portal build",            "0 errors", "—",      "npm run build --prefix client-portal"),
]
for i, (name, score, pct, src) in enumerate(eval_rows):
    y = 1.7 + i * 0.72
    bg = RGBColor(0x0A, 0x3D, 0x2B) if "100" in pct or "pass" in score.lower() or "0 " in score else CARD_BG
    rect(sl, 0.3, y, 12.73, 0.65, bg)
    txb(sl, name,  0.42, y+0.1, 5.5, 0.45, size=12, color=WHITE)
    txb(sl, score, 5.95, y+0.1, 1.3, 0.45, size=12, bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
    txb(sl, pct,   7.3,  y+0.1, 1.0, 0.45, size=12, bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
    txb(sl, src,   8.35, y+0.1, 4.55, 0.45, size=10, color=MID_GREY, italic=True)

# Quality thresholds
rect(sl, 0.3, 6.12, 12.73, 1.1, CARD_BG, ACCENT_GOLD)
txb(sl, "Definition of Done — Release Gate Thresholds",
    0.45, 6.17, 7.0, 0.38, size=13, bold=True, color=ACCENT_GOLD)
thresholds = [
    "Citation fidelity target: ≥ 95%",
    "Unsupported claim rate: ≤ 2%",
    "Weak-intent alert rate: < 70%",
    "Tier 3 Legal AI Done: grounding metadata on all outputs",
]
for i, t in enumerate(thresholds):
    txb(sl, "▸ " + t, 0.45 + i*3.18, 6.6, 3.1, 0.45, size=11, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 32 — SECTION: CONCLUSION
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
section_title(sl, "CONCLUSION & PERSPECTIVES", "09")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 33 — CONCLUSION
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Conclusion", "What was built and what was proven")

txb(sl, "Delivered a production-architecture legal AI platform covering the full lawyer workflow:",
    0.4, 1.2, 12.5, 0.45, size=15, color=WHITE)

achievements = [
    ("Backend Platform",    "FastAPI, multi-tenant RBAC, 15 API routers, full case/client/document lifecycle"),
    ("Document Intelligence","OCR → NER → PII redaction → chunking → FAISS indexing → hybrid retrieval"),
    ("AI Copilot",          "15+ intents, 12 specialized agents, grounded answers, citation verification"),
    ("Legal Search",        "Tunisia + Germany legal codes corpus, article-level retrieval, fallback taxonomy"),
    ("Voice Intake",        "Audio → transcription → consultation extraction → case linkage"),
    ("Client Portal",       "Public intake, OTP auth, case status, document upload — fully isolated"),
    ("Trust Layer",         "AI Insight: grounding type, confidence, legal_grounding, lawyer note"),
    ("Quality Gates",       "43/43 evals, 4/4 E2E, pytest suite, GitHub Actions release gate"),
]
for i, (title, desc) in enumerate(achievements):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.55
    y = 1.82 + row * 1.28
    rect(sl, x, y, 6.25, 1.15, CARD_BG)
    rect(sl, x, y, 1.85, 1.15, ACCENT_GOLD)
    txb(sl, title, x+0.08, y+0.3, 1.7, 0.5, size=12, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    txb(sl, desc,  x+1.98, y+0.18, 4.15, 0.8, size=11.5, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 34 — PERSPECTIVES
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
slide_header(sl, "Perspectives & Future Work", "Next steps toward production readiness")

short_term = [
    "Fine-tune embedding model on Tunisian legal corpus",
    "Article-level claim sentence mapping in UI",
    "Parser confidence scoring + disambiguation pass",
    "Alembic database migrations (replace schema_sync.py)",
    "LLM provider abstraction + exponential backoff",
]
long_term = [
    "Multi-jurisdiction expansion: France, Morocco, EU directives",
    "Custom fine-tuned legal LLM (instruction-tuned on case data)",
    "Real-time collaboration: multi-lawyer on same case",
    "Mobile app for lawyer field access",
    "SaaS multi-firm deployment with billing and SSO",
]

rect(sl, 0.3, 1.2, 6.1, 5.8, CARD_BG, ACCENT_GOLD)
rect(sl, 0.3, 1.2, 6.1, 0.5, ACCENT_GOLD)
txb(sl, "Short-Term  (next 3 months)", 0.42, 1.25, 5.8, 0.38, size=13, bold=True, color=DARK_NAVY)
for i, pt in enumerate(short_term):
    txb(sl, "▸ " + pt, 0.5, 1.85 + i*0.95, 5.7, 0.82, size=13, color=WHITE)

rect(sl, 6.7, 1.2, 6.35, 5.8, CARD_BG, ACCENT_GOLD)
rect(sl, 6.7, 1.2, 6.35, 0.5, ACCENT_GOLD)
txb(sl, "Long-Term  (6–18 months)", 6.82, 1.25, 6.1, 0.38, size=13, bold=True, color=DARK_NAVY)
for i, pt in enumerate(long_term):
    txb(sl, "▸ " + pt, 6.88, 1.85 + i*0.95, 6.0, 0.82, size=13, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 35 — THANK YOU
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_slide(sl)
rect(sl, 0, 0, 13.33, 0.5, ACCENT_GOLD)
rect(sl, 0, 7.0, 13.33, 0.5, ACCENT_GOLD)
rect(sl, 1.5, 1.8, 10.33, 3.9, CARD_BG)

txb(sl, "Thank You", 1.5, 2.0, 10.33, 1.3, size=64, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
gold_bar(sl, 3.38, 0.06)
txb(sl, "Questions & Discussion", 1.5, 3.5, 10.33, 0.7, size=24, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "Legal AI Platform  ·  PFE 2025–2026",
    1.5, 4.3, 10.33, 0.45, size=16, color=MID_GREY, align=PP_ALIGN.CENTER, italic=True)
txb(sl, "github.com/ahmedbm068/legal-ai-platform-",
    1.5, 4.85, 10.33, 0.4, size=14, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SAVE
# ════════════════════════════════════════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), "..", "docs", "Legal_AI_Platform_Soutenance.pptx")
prs.save(out)
print(f"Saved → {os.path.abspath(out)}")
print(f"Slides: {len(prs.slides)}")
